"""
MLB 승부예측 PPT 자동 생성기

실행: python auto_generate_ppt.py
출력: MLB_Daily_YYYYMMDD.pptx

predictions.csv를 읽어서 완성된 PPT를 자동 생성합니다.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import pandas as pd
from datetime import datetime, timedelta

CSV_FILE = 'predictions.csv'

# 색상
DARK_BG = RGBColor(15, 20, 25)
CARD_BG = RGBColor(26, 31, 46)
ACCENT = RGBColor(200, 241, 53)
RED = RGBColor(241, 53, 53)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(107, 114, 128)
LIGHT_GRAY = RGBColor(209, 213, 219)

def get_yesterday_stats():
    """어제 성적 계산"""
    df = pd.read_csv(CSV_FILE)
    yesterday = (datetime.now() - timedelta(days=1)).strftime('%Y-%m-%d')
    
    yesterday_games = df[df['date'] == yesterday]
    yesterday_correct = yesterday_games['correct'].dropna()
    
    if len(yesterday_correct) > 0:
        wins = int(yesterday_correct.sum())
        total = len(yesterday_correct)
        accuracy = (wins / total * 100) if total > 0 else 0
        return wins, total - wins, accuracy
    return 0, 0, 0.0

def get_week_stats():
    """주간 누적 성적"""
    df = pd.read_csv(CSV_FILE)
    week_ago = (datetime.now() - timedelta(days=7)).strftime('%Y-%m-%d')
    
    week_games = df[(df['date'] >= week_ago) & (df['correct'].notna())]
    if len(week_games) > 0:
        week_wins = int(week_games['correct'].sum())
        week_total = len(week_games)
        return week_wins, week_total - week_wins
    return 0, 0

def get_yesterday_picks():
    """어제 베스트픽 결과"""
    df = pd.read_csv(CSV_FILE)
    yesterday = (datetime.now() - timedelta(days=1)).strftime('%Y-%m-%d')
    
    yesterday_games = df[df['date'] == yesterday]
    
    picks = []
    for _, row in yesterday_games.iterrows():
        # 원정팀 체크
        if (row['model_away_prob'] >= 60 and 
            pd.notna(row.get('away_edge')) and 
            row['away_edge'] >= 5):
            picks.append({
                'team': row['away_team'],
                'prob': row['model_away_prob'],
                'correct': row['correct'],
                'away_score': row.get('away_score', '?'),
                'home_score': row.get('home_score', '?')
            })
        
        # 홈팀 체크
        if (row['model_home_prob'] >= 60 and 
            pd.notna(row.get('home_edge')) and 
            row['home_edge'] >= 5):
            picks.append({
                'team': row['home_team'],
                'prob': row['model_home_prob'],
                'correct': row['correct'],
                'away_score': row.get('away_score', '?'),
                'home_score': row.get('home_score', '?')
            })
    
    return picks[:3]  # TOP 3만

def get_today_picks():
    """오늘 베스트픽 TOP 3"""
    df = pd.read_csv(CSV_FILE)
    today = datetime.now().strftime('%Y-%m-%d')
    
    today_games = df[df['date'] == today]
    
    picks = []
    for _, row in today_games.iterrows():
        # 원정팀 체크
        if (row['model_away_prob'] >= 60 and 
            pd.notna(row.get('away_edge')) and 
            row['away_edge'] >= 5 and
            row['away_sp_confirmed']):
            picks.append({
                'team': row['away_team'],
                'opponent': row['home_team'],
                'prob': row['model_away_prob'],
                'edge': row['away_edge'],
                'sp': row['away_sp'],
                'opp_sp': row['home_sp'],
                'location': '원정'
            })
        
        # 홈팀 체크
        if (row['model_home_prob'] >= 60 and 
            pd.notna(row.get('home_edge')) and 
            row['home_edge'] >= 5 and
            row['home_sp_confirmed']):
            picks.append({
                'team': row['home_team'],
                'opponent': row['away_team'],
                'prob': row['model_home_prob'],
                'edge': row['home_edge'],
                'sp': row['home_sp'],
                'opp_sp': row['away_sp'],
                'location': '홈'
            })
    
    # Edge 높은 순 정렬
    picks.sort(key=lambda x: x['edge'], reverse=True)
    return picks[:3]

def create_daily_ppt():
    """PPT 자동 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    today_str = datetime.now().strftime('%Y년 %m월 %d일')
    yesterday_wins, yesterday_losses, yesterday_acc = get_yesterday_stats()
    week_wins, week_losses = get_week_stats()
    yesterday_picks = get_yesterday_picks()
    today_picks = get_today_picks()
    
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 슬라이드 1: 오프닝
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = DARK_BG
    
    # 상단 액센트 바
    accent_bar = slide1.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.05))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()
    
    # 타이틀
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.text = "MLB 승부예측"
    tf.paragraphs[0].font.size = Pt(72)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 날짜
    date_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    tf.text = today_str
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 성적 박스
    stats_box = slide1.shapes.add_shape(1, Inches(2), Inches(3.8), Inches(6), Inches(1.2))
    stats_box.fill.solid()
    stats_box.fill.fore_color.rgb = CARD_BG
    stats_box.line.fill.background()
    
    stats_text = slide1.shapes.add_textbox(Inches(2), Inches(4.1), Inches(6), Inches(0.6))
    tf = stats_text.text_frame
    p1 = tf.paragraphs[0]
    p1.text = f"어제 성적: {yesterday_wins}승 {yesterday_losses}패 ({yesterday_acc:.1f}%)"
    p1.font.size = Pt(20)
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = f"이번주 누적: {week_wins}승 {week_losses}패"
    p2.font.size = Pt(18)
    p2.font.color.rgb = ACCENT
    p2.alignment = PP_ALIGN.CENTER
    
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 슬라이드 2: 어제 결과
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    if len(yesterday_picks) > 0:
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        slide2.background.fill.solid()
        slide2.background.fill.fore_color.rgb = DARK_BG
        
        # 상단 액센트 바
        accent_bar = slide2.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.05))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = ACCENT
        accent_bar.line.fill.background()
        
        # 타이틀
        title = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
        tf = title.text_frame
        tf.text = "어제의 베스트픽 결과"
        tf.paragraphs[0].font.size = Pt(42)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        
        # 결과 카드
        y_pos = 1.5
        pick_wins = 0
        for pick in yesterday_picks:
            emoji = "✅" if pick['correct'] == 1 else "❌"
            result_text = f"승리!" if pick['correct'] == 1 else "패배"
            result_color = ACCENT if pick['correct'] == 1 else RED
            
            if pick['correct'] == 1:
                pick_wins += 1
            
            # 카드 배경
            card = slide2.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.9))
            card.fill.solid()
            card.fill.fore_color.rgb = CARD_BG
            card.line.fill.background()
            
            # 이모지
            emoji_box = slide2.shapes.add_textbox(Inches(0.7), Inches(y_pos+0.2), Inches(0.5), Inches(0.5))
            tf = emoji_box.text_frame
            tf.text = emoji
            tf.paragraphs[0].font.size = Pt(36)
            
            # 팀명 (박스 크기 늘림)
            team_box = slide2.shapes.add_textbox(Inches(1.4), Inches(y_pos+0.15), Inches(4.5), Inches(0.6))
            tf = team_box.text_frame
            tf.word_wrap = True
            tf.text = f"{pick['team']} ({pick['prob']:.1f}%)"
            tf.paragraphs[0].font.size = Pt(22)  # 24 → 22
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = WHITE
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # 결과 (박스 줄임)
            result_box = slide2.shapes.add_textbox(Inches(6), Inches(y_pos+0.15), Inches(3), Inches(0.6))
            tf = result_box.text_frame
            tf.text = f"→ {result_text}"
            tf.paragraphs[0].font.size = Pt(18)  # 20 → 18
            tf.paragraphs[0].font.color.rgb = result_color
            tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            y_pos += 1.0
        
        # 하단 요약
        pick_losses = len(yesterday_picks) - pick_wins
        pick_acc = (pick_wins / len(yesterday_picks) * 100) if len(yesterday_picks) > 0 else 0
        
        summary_box = slide2.shapes.add_shape(1, Inches(2.5), Inches(4.7), Inches(5), Inches(0.7))
        summary_box.fill.solid()
        summary_box.fill.fore_color.rgb = ACCENT
        summary_box.fill.fore_color.brightness = -0.5
        summary_box.line.fill.background()
        
        summary_text = slide2.shapes.add_textbox(Inches(2.5), Inches(4.7), Inches(5), Inches(0.7))
        tf = summary_text.text_frame
        tf.text = f"베스트픽: {pick_wins}승 {pick_losses}패 ({pick_acc:.1f}%)"
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 슬라이드 3~5: 오늘 베스트픽
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    medals = ["🥇", "🥈", "🥉"]
    medal_colors = [ACCENT, LIGHT_GRAY, LIGHT_GRAY]
    
    for idx, pick in enumerate(today_picks):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = DARK_BG
        
        # 상단 액센트 바
        accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.05))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = ACCENT
        accent_bar.line.fill.background()
        
        # 메달
        medal = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(0.8), Inches(0.8))
        tf = medal.text_frame
        tf.text = medals[idx]
        tf.paragraphs[0].font.size = Pt(48)
        
        # 순위
        rank = slide.shapes.add_textbox(Inches(1.4), Inches(0.5), Inches(8.1), Inches(0.6))
        tf = rank.text_frame
        tf.text = f"{idx+1}순위"
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = medal_colors[idx]
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 경기명 (폰트 줄이고 높이 늘림)
        game = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.8))
        tf = game.text_frame
        tf.text = f"{pick['team']} vs {pick['opponent']}"
        tf.paragraphs[0].font.size = Pt(28)  # 32 → 28
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.word_wrap = True
        
        # 예측 박스 (높이 늘림)
        pred_box = slide.shapes.add_shape(1, Inches(0.5), Inches(2.3), Inches(4.2), Inches(1.8))
        pred_box.fill.solid()
        pred_box.fill.fore_color.rgb = CARD_BG
        pred_box.line.fill.background()
        
        pred_text = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(3.8), Inches(1.4))
        tf = pred_text.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = "📊 예측"
        p1.font.size = Pt(18)  # 20 → 18
        p1.font.color.rgb = GRAY
        
        p2 = tf.add_paragraph()
        p2.text = f"{pick['team']} 승 {pick['prob']:.1f}%"
        p2.font.size = Pt(26)  # 28 → 26
        p2.font.bold = True
        p2.font.color.rgb = ACCENT
        p2.space_before = Pt(6)
        
        p3 = tf.add_paragraph()
        p3.text = f"Edge {pick['edge']:+.1f}%p"
        p3.font.size = Pt(18)  # 20 → 18
        p3.font.color.rgb = WHITE
        p3.space_before = Pt(3)
        
        # 선발 박스 (높이 늘림)
        sp_box = slide.shapes.add_shape(1, Inches(5.3), Inches(2.3), Inches(4.2), Inches(1.8))
        sp_box.fill.solid()
        sp_box.fill.fore_color.rgb = CARD_BG
        sp_box.line.fill.background()
        
        sp_text = slide.shapes.add_textbox(Inches(5.5), Inches(2.5), Inches(3.8), Inches(1.4))
        tf = sp_text.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = "⚾ 선발 매치업"
        p1.font.size = Pt(18)  # 20 → 18
        p1.font.color.rgb = GRAY
        
        p2 = tf.add_paragraph()
        p2.text = f"{pick['sp']}"
        p2.font.size = Pt(16)  # 18 → 16
        p2.font.color.rgb = WHITE
        p2.space_before = Pt(6)
        
        p3 = tf.add_paragraph()
        p3.text = "vs"
        p3.font.size = Pt(14)  # 16 → 14
        p3.font.color.rgb = GRAY
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(3)
        
        p4 = tf.add_paragraph()
        p4.text = f"{pick['opp_sp']}"
        p4.font.size = Pt(16)  # 18 → 16
        p4.font.color.rgb = WHITE
        p4.space_before = Pt(3)
        
        # 주목 포인트 (위치 조정)
        points_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(0.4))
        tf = points_title.text_frame
        tf.text = "💡 주목 포인트"
        tf.paragraphs[0].font.size = Pt(20)  # 22 → 20
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        
        placeholder = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.7))
        tf = placeholder.text_frame
        tf.word_wrap = True
        tf.text = "[여기에 주목 포인트 입력]\n예: 최근 5일 타선 폭발, 불펜 격차 등"
        tf.paragraphs[0].font.size = Pt(14)  # 16 → 14
        tf.paragraphs[0].font.color.rgb = GRAY
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 마지막 슬라이드: 요약
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide_final = prs.slides.add_slide(prs.slide_layouts[6])
    slide_final.background.fill.solid()
    slide_final.background.fill.fore_color.rgb = DARK_BG
    
    accent_bar = slide_final.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.05))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()
    
    title = slide_final.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.6))
    tf = title.text_frame
    tf.text = "오늘의 베스트픽 요약"
    tf.paragraphs[0].font.size = Pt(42)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    summary_box = slide_final.shapes.add_shape(1, Inches(2), Inches(1.8), Inches(6), Inches(2.2))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = CARD_BG
    summary_box.line.fill.background()
    
    summary_text = slide_final.shapes.add_textbox(Inches(2.5), Inches(2.1), Inches(5), Inches(1.6))
    tf = summary_text.text_frame
    
    for idx, pick in enumerate(today_picks):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"{medals[idx]} {pick['team']} ({pick['prob']:.1f}%, Edge {pick['edge']:+.1f}%p)"
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        if idx > 0:
            p.space_before = Pt(10)
    
    info = slide_final.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(9), Inches(0.4))
    tf = info.text_frame
    tf.text = "자세한 스탯은 설명란에서 확인!"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = GRAY
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    cta_box = slide_final.shapes.add_shape(1, Inches(3), Inches(4.9), Inches(4), Inches(0.5))
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = ACCENT
    cta_box.line.fill.background()
    
    cta_text = slide_final.shapes.add_textbox(Inches(3), Inches(4.9), Inches(4), Inches(0.5))
    tf = cta_text.text_frame
    tf.text = "구독 & 좋아요 부탁드립니다!"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK_BG
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 저장
    filename = f"MLB_Daily_{datetime.now().strftime('%Y%m%d')}.pptx"
    prs.save(filename)
    
    print(f"✅ {filename} 생성 완료!")
    print(f"   슬라이드: {len(prs.slides)}개")
    print(f"   어제 성적: {yesterday_wins}승 {yesterday_losses}패")
    print(f"   오늘 베스트픽: {len(today_picks)}개")
    print(f"\n⚠️  주목 포인트는 수동으로 작성해주세요!")

if __name__ == "__main__":
    try:
        create_daily_ppt()
    except FileNotFoundError:
        print("❌ predictions.csv 파일을 찾을 수 없습니다.")
        print("   predict.py daily를 먼저 실행하세요.")
    except Exception as e:
        print(f"❌ 에러 발생: {e}")
