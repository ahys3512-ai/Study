"""
MLB 승부예측 PPT 템플릿 생성기

실행: python create_mlb_template.py
출력: MLB_Template.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# 색상 팔레트
DARK_BG = RGBColor(15, 20, 25)       # 0F1419
CARD_BG = RGBColor(26, 31, 46)       # 1A1F2E
ACCENT = RGBColor(200, 241, 53)      # C8F135
RED = RGBColor(241, 53, 53)          # F13535
WHITE = RGBColor(255, 255, 255)      # FFFFFF
GRAY = RGBColor(107, 114, 128)       # 6B7280
LIGHT_GRAY = RGBColor(209, 213, 219) # D1D5DB

def create_mlb_template():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 슬라이드 1: 오프닝
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = DARK_BG
    
    # 상단 액센트 바
    accent_bar = slide1.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.05))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()
    
    # 메인 타이틀
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.text = "MLB 승부예측"
    tf.paragraphs[0].font.size = Pt(72)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 날짜
    date_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    tf.text = "2026년 5월 7일"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 하단 박스
    stats_box = slide1.shapes.add_shape(1, Inches(2), Inches(3.8), Inches(6), Inches(1.2))
    stats_box.fill.solid()
    stats_box.fill.fore_color.rgb = CARD_BG
    stats_box.line.fill.background()
    
    stats_text = slide1.shapes.add_textbox(Inches(2), Inches(4.1), Inches(6), Inches(0.6))
    tf = stats_text.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "어제 성적: 8승 7패 (53.3%)"
    p1.font.size = Pt(20)
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "이번주 누적: 8승 7패"
    p2.font.size = Pt(18)
    p2.font.color.rgb = ACCENT
    p2.alignment = PP_ALIGN.CENTER
    
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 슬라이드 2: 어제 결과
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = DARK_BG
    
    # 상단 액센트 바
    accent_bar = slide2.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.05))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()
    
    # 타이틀
    title = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.6))
    tf = title.text_frame
    tf.text = "어제의 베스트픽 결과"
    tf.paragraphs[0].font.size = Pt(42)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    
    # 결과 카드들
    results = [
        ("✅", "Atlanta Braves (67.5%)", "→ 3-2 승리!", ACCENT),
        ("✅", "Minnesota Twins (68.3%)", "→ 11-3 대승!", ACCENT),
        ("❌", "Chicago White Sox (61.2%)", "→ 3-4 석패", RED)
    ]
    
    y_pos = 1.5
    for emoji, team, result, result_color in results:
        # 카드 배경
        card = slide2.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.9))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.fill.background()
        
        # 이모지
        emoji_box = slide2.shapes.add_textbox(Inches(0.7), Inches(y_pos+0.2), Inches(0.5), Inches(0.5))
        tf = emoji_box.text_frame
        tf.text = emoji
        tf.paragraphs[0].font.size = Pt(36)
        
        # 팀명
        team_box = slide2.shapes.add_textbox(Inches(1.4), Inches(y_pos+0.2), Inches(4), Inches(0.5))
        tf = team_box.text_frame
        tf.text = team
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 결과
        result_box = slide2.shapes.add_textbox(Inches(5.5), Inches(y_pos+0.2), Inches(3.5), Inches(0.5))
        tf = result_box.text_frame
        tf.text = result
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = result_color
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        y_pos += 1.0
    
    # 하단 요약
    summary_box = slide2.shapes.add_shape(1, Inches(2.5), Inches(4.7), Inches(5), Inches(0.7))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = ACCENT
    summary_box.fill.fore_color.brightness = -0.5
    summary_box.line.fill.background()
    
    summary_text = slide2.shapes.add_textbox(Inches(2.5), Inches(4.7), Inches(5), Inches(0.7))
    tf = summary_text.text_frame
    tf.text = "베스트픽: 2승 1패 (66.7%)"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 슬라이드 3: 1순위 베스트픽
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    slide3.background.fill.solid()
    slide3.background.fill.fore_color.rgb = DARK_BG
    
    # 상단 액센트 바
    accent_bar = slide3.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.05))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()
    
    # 메달 + 순위
    medal = slide3.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(0.8), Inches(0.8))
    tf = medal.text_frame
    tf.text = "🥇"
    tf.paragraphs[0].font.size = Pt(48)
    
    rank = slide3.shapes.add_textbox(Inches(1.4), Inches(0.5), Inches(8.1), Inches(0.6))
    tf = rank.text_frame
    tf.text = "1순위"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 경기명
    game = slide3.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.6))
    tf = game.text_frame
    tf.text = "New York Yankees vs Boston Red Sox"
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    
    # 왼쪽 박스: 예측
    pred_box = slide3.shapes.add_shape(1, Inches(0.5), Inches(2.2), Inches(4.2), Inches(1.5))
    pred_box.fill.solid()
    pred_box.fill.fore_color.rgb = CARD_BG
    pred_box.line.fill.background()
    
    pred_text = slide3.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(3.8), Inches(1.1))
    tf = pred_text.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "📊 예측"
    p1.font.size = Pt(20)
    p1.font.color.rgb = GRAY
    
    p2 = tf.add_paragraph()
    p2.text = "Yankees 승 65.3%"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT
    p2.space_before = Pt(6)
    
    p3 = tf.add_paragraph()
    p3.text = "Edge +8.2%p"
    p3.font.size = Pt(20)
    p3.font.color.rgb = WHITE
    p3.space_before = Pt(3)
    
    # 오른쪽 박스: 선발
    sp_box = slide3.shapes.add_shape(1, Inches(5.3), Inches(2.2), Inches(4.2), Inches(1.5))
    sp_box.fill.solid()
    sp_box.fill.fore_color.rgb = CARD_BG
    sp_box.line.fill.background()
    
    sp_text = slide3.shapes.add_textbox(Inches(5.5), Inches(2.4), Inches(3.8), Inches(1.1))
    tf = sp_text.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "⚾ 선발 매치업"
    p1.font.size = Pt(20)
    p1.font.color.rgb = GRAY
    
    p2 = tf.add_paragraph()
    p2.text = "Gerrit Cole (ERA 2.89)"
    p2.font.size = Pt(18)
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(6)
    
    p3 = tf.add_paragraph()
    p3.text = "vs"
    p3.font.size = Pt(16)
    p3.font.color.rgb = GRAY
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(3)
    
    p4 = tf.add_paragraph()
    p4.text = "Brayan Bello (ERA 4.12)"
    p4.font.size = Pt(18)
    p4.font.color.rgb = WHITE
    p4.space_before = Pt(3)
    
    # 하단: 주목 포인트
    points_title = slide3.shapes.add_textbox(Inches(0.5), Inches(3.9), Inches(9), Inches(0.4))
    tf = points_title.text_frame
    tf.text = "💡 주목 포인트"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    
    points = [
        "✓ Yankees 최근 5일 OPS 0.856 (폭발 중)",
        "✓ 불펜 ERA 격차: 양키스 2.12 vs 보스턴 3.98",
        "✓ 홈 경기 이점"
    ]
    
    y_pos = 4.4
    for point in points:
        point_box = slide3.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.3))
        point_box.fill.solid()
        point_box.fill.fore_color.rgb = CARD_BG
        point_box.line.fill.background()
        
        point_text = slide3.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(8.6), Inches(0.3))
        tf = point_text.text_frame
        tf.text = point
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        y_pos += 0.35
    
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 슬라이드 4, 5: 2순위, 3순위 (간단 버전)
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    for rank_num, medal_emoji, game_title in [
        (2, "🥈", "Los Angeles Dodgers vs San Francisco Giants"),
        (3, "🥉", "Houston Astros vs Seattle Mariners")
    ]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = DARK_BG
        
        # 상단 액센트 바
        accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.05))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = ACCENT
        accent_bar.line.fill.background()
        
        # 메달
        medal = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(0.8), Inches(0.8))
        tf = medal.text_frame
        tf.text = medal_emoji
        tf.paragraphs[0].font.size = Pt(48)
        
        # 순위
        rank = slide.shapes.add_textbox(Inches(1.4), Inches(0.5), Inches(8.1), Inches(0.6))
        tf = rank.text_frame
        tf.text = f"{rank_num}순위"
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 경기명
        game = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.6))
        tf = game.text_frame
        tf.text = game_title
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        
        # 플레이스홀더
        placeholder = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2.5))
        tf = placeholder.text_frame
        tf.text = f"[여기에 {rank_num}순위 경기 데이터 입력]\n\n예측, 선발, 주목포인트 등"
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = GRAY
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 슬라이드 6: 마무리
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    slide6.background.fill.solid()
    slide6.background.fill.fore_color.rgb = DARK_BG
    
    # 상단 액센트 바
    accent_bar = slide6.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.05))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()
    
    # 타이틀
    title = slide6.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.6))
    tf = title.text_frame
    tf.text = "오늘의 베스트픽 요약"
    tf.paragraphs[0].font.size = Pt(42)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 요약 박스
    summary_box = slide6.shapes.add_shape(1, Inches(2), Inches(1.8), Inches(6), Inches(2.2))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = CARD_BG
    summary_box.line.fill.background()
    
    summary_text = slide6.shapes.add_textbox(Inches(2.5), Inches(2.1), Inches(5), Inches(1.6))
    tf = summary_text.text_frame
    
    p1 = tf.paragraphs[0]
    p1.text = "🥇 Yankees (65.3%, Edge +8.2%p)"
    p1.font.size = Pt(22)
    p1.font.color.rgb = WHITE
    p1.space_after = Pt(10)
    
    p2 = tf.add_paragraph()
    p2.text = "🥈 Dodgers (61.5%, Edge +7.1%p)"
    p2.font.size = Pt(22)
    p2.font.color.rgb = WHITE
    p2.space_after = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = "🥉 Astros (60.2%, Edge +5.8%p)"
    p3.font.size = Pt(22)
    p3.font.color.rgb = WHITE
    
    # 하단 안내
    info = slide6.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(9), Inches(0.4))
    tf = info.text_frame
    tf.text = "자세한 스탯은 설명란에서 확인!"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = GRAY
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # CTA 버튼
    cta_box = slide6.shapes.add_shape(1, Inches(3), Inches(4.9), Inches(4), Inches(0.5))
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = ACCENT
    cta_box.line.fill.background()
    
    cta_text = slide6.shapes.add_textbox(Inches(3), Inches(4.9), Inches(4), Inches(0.5))
    tf = cta_text.text_frame
    tf.text = "구독 & 좋아요 부탁드립니다!"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK_BG
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 저장
    prs.save('MLB_Template.pptx')
    print("✅ MLB_Template.pptx 생성 완료!")
    print("   6개 슬라이드: 오프닝, 어제결과, 1순위, 2순위, 3순위, 마무리")

if __name__ == "__main__":
    create_mlb_template()
